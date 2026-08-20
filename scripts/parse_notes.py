#!/usr/bin/env python3
"""
study-notes-summarizer — unified free parser (cross-platform)
==============================================================
Parse a folder of mixed-format study notes into Markdown text, written to the parsed/
directory, and record a processed.log to avoid re-processing. Uses only free local
tools throughout — no paid API required.

Supported formats and their free engines:
  .docx            -> python-docx
  .doc (legacy)    -> cross-platform fallback chain:
                       1) macOS built-in textutil
                       2) Windows Microsoft Word + pywin32 (COM)
                       3) LibreOffice (soffice/libreoffice) convert to docx, then parse
                       4) Linux antiword
                       if none available, print a clear install hint (or Save As .docx manually)
  .pdf (text)      -> pdfplumber
  .pdf (scanned)   -> PyMuPDF render + tesseract OCR (chi_sim+deu+eng, language pack auto-download)
  .pdf (watermark) -> text layer that is mostly repeated header/watermark lines
                      (real content in page images) is auto-detected and the
                      whole file is re-OCR'd; requires the OCR engine above.
  .png/.jpg/...    -> tesseract OCR (chi_sim+deu+eng, optionally Save As PDF via img2pdf)
  .pptx            -> python-pptx
  .ppt (legacy)    -> cross-platform fallback chain: textutil -> LibreOffice to pptx -> hint
  .xlsx            -> openpyxl (each sheet rendered as a Markdown table)
  .mp3/.m4a/.wav/.flac/.ogg/.aac and video containers (.mp4/.mov/.webm/.m4v)
                  -> Whisper speech-to-text (faster-whisper, local + free, language
                     auto-detected; model auto-downloaded on first use). Audio decoding
                     uses PyAV, which bundles its own ffmpeg — no system install needed.
  .zip / .rar      -> auto-extract and recurse into contents (PDF/DOCX/.../audio inside)

When an engine is missing, the script prints a clear install hint and skips that file
(not written to processed.log, so it auto-retries next run once the engine is installed)
— it never fails silently.

Usage:
  python3 parse_notes.py --source "/path/to/notes" --out "/path/to/parsed" [--log processed.log] [--force]
"""
import argparse
import os
import sys
import subprocess
import shutil
from pathlib import Path

# ---------- tool availability detection ----------
def has(cmd):
    return shutil.which(cmd) is not None

def py_has(mod):
    try:
        __import__(mod)
        return True
    except Exception:
        return False

def find_executable(name, extra_dirs=None):
    """Prefer PATH; otherwise search common install dirs / the skill's bundled bin / env vars
    (Windows installers often skip PATH). Returns the full path or None."""
    # 1) explicit override (lets the user drop in a portable binary)
    env = os.environ.get("NOTES_TESSERACT_BIN") if name == "tesseract" else None
    if env and os.path.isfile(env):
        return env
    # 2) PATH
    p = shutil.which(name)
    if p:
        return p
    import glob as _glob
    prog = os.environ.get("ProgramFiles", r"C:\Program Files")
    prog86 = os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")
    local = os.environ.get("LOCALAPPDATA", "")
    skill_bin = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "bin")
    candidates = list(extra_dirs or [])
    candidates += [
        # skill-bundled portable copy (fully zero-manual)
        os.path.join(skill_bin, "tesseract", name + ".exe"),
        os.path.join(skill_bin, name + ".exe"),
        # common Windows installer locations that often skip PATH
        os.path.join(prog, "Tesseract-OCR", name + ".exe"),
        os.path.join(prog86, "Tesseract-OCR", name + ".exe"),
        os.path.join(prog, "Poppler", "bin", name + ".exe"),
        os.path.join(prog86, "Poppler", "bin", name + ".exe"),
        os.path.join(prog, "poppler-*", "bin", name + ".exe"),
        os.path.join(local, "Programs", "Poppler", "bin", name + ".exe"),
    ]
    for c in candidates:
        for m in (_glob.glob(c) or [c]):
            if os.path.isfile(m):
                return m
    return None


def _safe_unlink(p):
    """Delete a temp file; ignore failures (some sandboxes/recycle-bin hooks block unlink and must not break the main flow)."""
    try:
        Path(p).unlink(missing_ok=True)
    except Exception:
        pass


# ---------- zero-manual: dependency bootstrap & tesseract auto-install ----------
def _pip_install(pkgs):
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "--quiet",
                       "--disable-pip-version-check", *pkgs],
                      check=True, capture_output=True)
        return True
    except Exception as e:
        print(f"⚠️ Auto-install of Python dependencies failed (run `pip install {' '.join(pkgs)}` manually): {e}",
              file=sys.stderr)
        return False


def _bootstrap_deps():
    """On first run, auto-install missing Python dependencies so the end user never runs pip manually.
    Set NOTES_SKIP_DEP_INSTALL=1 to disable (e.g. offline environments)."""
    if os.environ.get("NOTES_SKIP_DEP_INSTALL") == "1":
        return
    required = {
        "docx": "python-docx", "pptx": "python-pptx", "pdfplumber": "pdfplumber",
        "pytesseract": "pytesseract", "PIL": "pillow", "fitz": "pymupdf",
        "img2pdf": "img2pdf", "rarfile": "rarfile", "openpyxl": "openpyxl",
        "faster_whisper": "faster-whisper",
    }
    # pywin32 is Windows-only (.doc via Word COM); macOS uses native textutil, Linux uses antiword/LibreOffice.
    if sys.platform.startswith("win"):
        required["win32com"] = "pywin32"
    optional = {"pdf2image": "pdf2image"}
    missing = [pip for mod, pip in required.items() if not py_has(mod)]
    if missing:
        print(f"📦 First run: auto-installing Python dependencies: {', '.join(missing)}", file=sys.stderr)
        _pip_install(missing)
    miss_opt = [pip for mod, pip in optional.items() if not py_has(mod)]
    if miss_opt:
        _pip_install(miss_opt)  # best-effort; failure does not block the main flow
    # Windows: after pywin32 installs, run postinstall once or the first import win32com fails
    if sys.platform.startswith("win") and "pywin32" in missing:
        _try_pywin32_postinstall()


def _try_pywin32_postinstall():
    """pywin32 in a venv usually needs a one-time postinstall before import win32com works."""
    try:
        script = os.path.join(os.path.dirname(sys.executable), "pywin32_postinstall.py")
        if os.path.isfile(script):
            subprocess.run([sys.executable, script, "-install"],
                          check=False, capture_output=True)
    except Exception:
        pass


def _run_installer(cmd, args):
    try:
        subprocess.run([cmd, *args], check=True)
        return True
    except Exception as e:
        print(f"⚠️ Auto-install of {cmd} failed; install tesseract manually per the prompt: {e}", file=sys.stderr)
        return False


def _ensure_tesseract():
    """When tesseract is missing, auto-call the platform's native package manager to install it
    (zero-manual), then re-detect. Windows→winget; macOS→brew. Linux (apt) needs sudo, so we print
    a manual hint instead. Set NOTES_SKIP_TESSERACT_INSTALL=1 to disable."""
    global TESSERACT_BIN
    if TESSERACT_BIN:
        return
    if os.environ.get("NOTES_SKIP_TESSERACT_INSTALL") == "1":
        return
    print("🔍 tesseract not found, attempting auto-install (zero-manual)...", file=sys.stderr)
    ok = False
    if sys.platform.startswith("win"):
        ok = _run_installer("winget", ["install", "--accept-package-agreements",
                                       "--accept-source-agreements", "-e",
                                       "--id", "UB-Mannheim.TesseractOCR"])
    elif sys.platform == "darwin":
        ok = _run_installer("brew", ["install", "tesseract"])
    if ok:
        TESSERACT_BIN = find_executable("tesseract")
        if TESSERACT_BIN:
            print(f"✅ tesseract ready: {TESSERACT_BIN}", file=sys.stderr)
        else:
            print("⚠️ Install ran but tesseract could not be located; restart the terminal or add it to PATH manually.",
                  file=sys.stderr)


_bootstrap_deps()   # zero-manual: auto-install any missing Python dependencies

TEXTUTIL      = has("textutil")
PDFPLUMBER   = py_has("pdfplumber")
PYTHON_DOCX  = py_has("docx")
PYTHON_PPTX  = py_has("pptx")
IMG2PDF      = py_has("img2pdf")
PYTESSERACT  = py_has("pytesseract")
TESSERACT_BIN = find_executable("tesseract")
PDFTOPPM_BIN  = find_executable("pdftoppm")   # part of poppler, converts scanned PDF to images
PDF2IMAGE    = py_has("pdf2image")       # pdfplumber scanned-OCR helper (poppler path)
PIL_OK       = py_has("PIL")             # image OCR
PYMUPDF      = py_has("fitz")            # PyMuPDF, scanned-PDF rendering (replaces poppler, pure pip)
WIN32COM     = py_has("win32com")        # Windows Word COM parses .doc
OPENPYXL     = py_has("openpyxl")        # openpyxl parses .xlsx
SOFFICE      = has("soffice") or has("libreoffice")
ANTIWORD     = has("antiword")           # common Linux .doc text extractor

_ensure_tesseract()   # zero-manual: auto-install tesseract per platform when missing, then re-detect

# OCR availability
OCR_IMG_OK  = TESSERACT_BIN and PYTESSERACT and PIL_OK          # image OCR
OCR_PDF_OK  = OCR_IMG_OK and (PYMUPDF or (PDFTOPPM_BIN and PDF2IMAGE))  # scanned-PDF OCR (PyMuPDF preferred)

# Speech-to-text (Whisper) availability — faster-whisper decodes audio via PyAV,
# which bundles its own ffmpeg, so no system ffmpeg binary is required.
WHISPER_OK  = py_has("faster_whisper")       # local, free speech-to-text
AUDIO_OK    = WHISPER_OK                      # PyAV (ffmpeg) ships with faster-whisper

# ---------- dispatch table ----------
EXT_DISPATCH = {
    ".docx": "parse_docx",
    ".doc": "parse_doc",
    ".pdf": "parse_pdf",
    ".pptx": "parse_pptx",
    ".ppt": "parse_ppt",
    ".xlsx": "parse_xlsx",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"}
ARCHIVE_EXTS = {".zip", ".rar"}   # auto-extract and recurse into contents
AUDIO_EXTS = {".mp3", ".m4a", ".wav", ".flac", ".ogg", ".aac", ".wma",
              ".mp4", ".mov", ".webm", ".m4v"}   # speech-to-text via Whisper

SUPPORTED = set(EXT_DISPATCH) | IMAGE_EXTS | ARCHIVE_EXTS | AUDIO_EXTS


# ---------- per-format parsers (lazy-loaded; prompt if missing) ----------
def parse_docx(path: Path):
    import docx
    d = docx.Document(str(path))
    lines = []
    for p in d.paragraphs:
        if p.text.strip():
            lines.append(p.text)
    for ti, table in enumerate(d.tables, 1):
        lines.append(f"\n[Table {ti}]")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            lines.append(" | ".join(cells))
    return "\n".join(lines), "python-docx"


def parse_doc(path: Path, out_dir: Path):
    """Cross-platform .doc parsing.
    Preferred path: convert .doc -> .docx with the best available engine (keeps tables),
    then read via python-docx. Falls back to text-only extraction when docx conversion
    is unavailable or fails (macOS textutil -convert txt, Linux antiword)."""
    # 1) macOS built-in textutil: .docx conversion first, plain-text fallback
    if TEXTUTIL:
        docx_text = _doc_via_textutil_docx(path, out_dir)
        if docx_text is not None:
            return docx_text, "macOS textutil (docx)"
        return _doc_via_textutil(path, out_dir), "macOS textutil (txt)"
    # 2) Windows Microsoft Word COM -> .docx -> python-docx
    if sys.platform.startswith("win") and WIN32COM:
        try:
            return _doc_via_winword(path, out_dir), "Windows Word COM"
        except Exception as e:
            raise RuntimeError(
                f"Windows Word COM failed to parse .doc: {e}. "
                f"Please ensure Microsoft Word is installed and run `pip install pywin32`."
            )
    # 3) LibreOffice -> .docx -> python-docx
    if SOFFICE:
        return _doc_via_soffice(path, out_dir, "docx"), "LibreOffice"
    # 4) Linux antiword (text-only)
    if ANTIWORD:
        return _doc_via_antiword(path), "antiword"
    raise RuntimeError(
        "Parsing .doc requires one of the following (all checked by priority, none found):\n"
        "  1) macOS built-in textutil;\n"
        "  2) Windows with Microsoft Word installed + `pip install pywin32`;\n"
        "  3) LibreOffice (soffice / libreoffice);\n"
        "  4) Linux antiword.\n"
        "Please install one of these, or 'Save As .docx' in Word/WPS before feeding the file."
    )


def _doc_via_textutil(path: Path, out_dir: Path):
    tmp = out_dir / f"._{path.stem}.txt"
    subprocess.run(["textutil", "-convert", "txt", "-output", str(tmp), str(path)],
                   check=True, capture_output=True)
    text = tmp.read_text(encoding="utf-8", errors="ignore")
    _safe_unlink(tmp)
    return text


def _doc_via_textutil_docx(path: Path, out_dir: Path):
    """macOS: convert .doc -> .docx with the native textutil (preserves tables),
    then read with python-docx. Returns the extracted text, or None if conversion
    or reading fails, so the caller can fall back to plain-text extraction."""
    tmp_docx = out_dir / f"._{path.stem}_textutil.docx"
    try:
        subprocess.run(["textutil", "-convert", "docx", "-output", str(tmp_docx), str(path)],
                       check=True, capture_output=True)
        if not tmp_docx.exists():
            return None
        return parse_docx(tmp_docx)[0]
    except Exception:
        return None
    finally:
        _safe_unlink(tmp_docx)


def _doc_via_winword(path: Path, out_dir: Path):
    """Use Windows-native Microsoft Word to save .doc as .docx, then read with python-docx (preserves tables)."""
    import pythoncom
    import win32com.client
    pythoncom.CoInitialize()
    word = None
    converted = out_dir / f"._{path.stem}_converted.docx"
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(path))
        # wdFormatXMLDocument = 12
        doc.SaveAs(str(converted), FileFormat=12)
        doc.Close(False)
    finally:
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()
    try:
        text, _ = parse_docx(converted)
    finally:
        _safe_unlink(converted)
    return text


def _doc_via_soffice(path: Path, out_dir: Path, target_ext: str):
    soffice_bin = find_executable("soffice") or find_executable("libreoffice") or "soffice"
    converted = out_dir / f"{path.stem}.{target_ext}"
    subprocess.run(
        [soffice_bin, "--headless", "--convert-to", target_ext,
         "--outdir", str(out_dir), str(path)],
        check=True, capture_output=True,
    )
    if not converted.exists():
        raise RuntimeError("LibreOffice did not produce the expected file: " + str(converted))
    try:
        if target_ext == "docx":
            text, _ = parse_docx(converted)
        else:
            text, _ = parse_pptx(converted)
    finally:
        _safe_unlink(converted)
    return text


def _doc_via_antiword(path: Path):
    r = subprocess.run(["antiword", "-t", str(path)],
                       capture_output=True, text=True)
    if r.returncode != 0:
        raise RuntimeError("antiword conversion failed: " + (r.stderr or r.stdout))
    return r.stdout


def _is_watermark_overlay(page_text: dict) -> bool:
    """Detect PDFs whose text layer is dominated by repeated header/watermark
    lines while the real content lives in page images. Typical of slide-deck
    PDFs that stamp a course-name/logo watermark on every page — pdfplumber
    extracts plenty of chars, but they are pure noise and the actual content
    is silently lost. Returns True when repeated lines dominate (>40% of all
    lines appear 3+ times) or almost no unique content lines remain."""
    lines = []
    for t in page_text.values():
        lines.extend(l.strip() for l in t.splitlines() if l.strip())
    if not lines:
        return False
    counts = {}
    for l in lines:
        counts[l] = counts.get(l, 0) + 1
    repeated = sum(n for n in counts.values() if n >= 3)
    unique_content = sum(1 for l, n in counts.items() if n < 3 and len(l) >= 4)
    return (repeated / len(lines)) > 0.4 or unique_content < 15


def parse_pdf(path: Path, out_dir: Path):
    """Extract text page-by-page (Bug 4/7): a page is OCR'd only when it has
    little native text (< 30 chars), so mixed PDFs keep their text pages and
    get scanned pages OCR'd instead of being silently dropped; short text PDFs
    are never mis-detected as scans.
    Watermark-overlay fix: when every page has 30+ chars but the text layer is
    mostly repeated watermarks, the whole file is re-OCR'd because trusting the
    noise would drop the real (image-based) content."""
    import pdfplumber
    page_text = {}
    need_ocr = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            t = page.extract_text() or ""
            page_text[i] = t
            if len(t.strip()) < 30:
                need_ocr.append(i)

    # Watermark-overlay fix: all pages have text, but it is mostly repeated
    # header/watermark noise while the real content sits in page images.
    if not need_ocr and len(page_text) >= 3 and _is_watermark_overlay(page_text):
        print(f"ℹ️ `{path.name}`: text layer looks like repeated watermarks "
              f"(real content is in page images) — re-OCR'ing the whole file.",
              file=sys.stderr)
        need_ocr = sorted(page_text)

    ocr_result = {}
    all_scanned = len(page_text) > 0 and len(need_ocr) == len(page_text)
    if need_ocr:
        if OCR_PDF_OK:
            try:
                ocr_result = _ocr_pdf(path, out_dir, page_indices=need_ocr)
            except Exception as e:
                print(f"⚠️ OCR failed for some pages of `{path.name}`: {e}", file=sys.stderr)
        elif all_scanned:
            # The whole file is a scan and OCR is unavailable -> must tell the user
            missing = []
            if not TESSERACT_BIN: missing.append("tesseract")
            if not PYTESSERACT:   missing.append("pytesseract (Python lib)")
            if not (PYMUPDF or (PDFTOPPM_BIN and PDF2IMAGE)): missing.append("scanned-PDF renderer (PyMuPDF or poppler)")
            if not PIL_OK:        missing.append("Pillow (PIL)")
            raise RuntimeError(
                f"PDF `{path.name}` looks like a scan or image-based slide deck "
                f"(little real text extractable) and needs OCR, but is missing: {', '.join(missing)}."
                f"\n  Recommended (pure pip, no system binary): `pip install pytesseract pymupdf pillow`."
                f"\n  Or the poppler route: macOS `brew install tesseract poppler`;"
                f"Linux `apt install tesseract-ocr poppler-utils`;"
                f"then `pip install pytesseract pdf2image pillow`."
                f"\n  tesseract itself: Windows `winget install UB-Mannheim.TesseractOCR`;"
                f"macOS `brew install tesseract`; Linux `apt install tesseract-ocr`."
            )
        else:
            # Mixed PDF: some pages need OCR but the engine is missing — mark those
            # pages instead of crashing the whole run.
            missing = []
            if not TESSERACT_BIN: missing.append("tesseract")
            if not PYTESSERACT:   missing.append("pytesseract (Python lib)")
            if not (PYMUPDF or (PDFTOPPM_BIN and PDF2IMAGE)): missing.append("scanned-PDF renderer (PyMuPDF or poppler)")
            if not PIL_OK:        missing.append("Pillow (PIL)")
            print(f"⚠️ `{path.name}` has page(s) with little text needing OCR but is missing: {', '.join(missing)}.",
                  file=sys.stderr)

    bodies = []
    ocr_pages = 0
    for i in sorted(page_text):
        t = page_text[i]
        ocr = ocr_result.get(i)
        if ocr is not None and ocr.strip():
            bodies.append(f"[Page {i}]\n{ocr}")
            ocr_pages += 1
        elif t.strip():
            bodies.append(f"[Page {i}]\n{t}")
        else:
            bodies.append(f"[Page {i}]\n(scan page — OCR unavailable or produced no text)")
    joined = "\n\n".join(bodies)
    label = "pdfplumber" + (f" (+OCR x{ocr_pages})" if ocr_pages else "")
    return joined, label


# ---------- OCR language-pack resolution / auto-download ----------
def _writable(d):
    """Reliable writability test: os.access(d, os.W_OK) is unreliable on Windows
    (often reports Program Files as writable when it is not), so we actually try
    to create a temp file. Whether we can *delete* the probe is irrelevant to
    writability, so cleanup failures are ignored."""
    if not d or not os.path.isdir(d):
        return False
    probe = os.path.join(d, f".notes_write_test_{os.getpid()}")
    try:
        with open(probe, "w") as f:
            f.write("")
    except Exception:
        return False
    try:
        os.remove(probe)
    except Exception:
        pass
    return True


def _default_tessdata_dir():
    """Locate the tesseract language-pack directory (tessdata).

    Order of preference:
      1. `tesseract --print-tessdata-dir` output (authoritative when supported)
      2. <tesseract dir>/tessdata  (portable/Windows layouts)
      3. common Homebrew layouts: <prefix>/share/tessdata for Intel (/usr/local)
         and Apple Silicon (/opt/homebrew), where the binary lives in bin/
    Returns the dir as a str, or None."""
    if TESSERACT_BIN:
        try:
            out = subprocess.run(
                [TESSERACT_BIN, "--print-tessdata-dir"],
                capture_output=True, text=True, timeout=15,
            )
            d = (out.stdout or "").strip()
            if d and os.path.isdir(d):
                return d
        except Exception:
            pass
        bdir = os.path.dirname(TESSERACT_BIN)
        for cand in (
            os.path.join(bdir, "tessdata"),
            os.path.join(bdir, "..", "share", "tessdata"),
            "/usr/local/share/tessdata",
            "/opt/homebrew/share/tessdata",
        ):
            cand = os.path.abspath(cand)
            if os.path.isdir(cand):
                return cand
    return None


def _ensure_tessdata_dir():
    """Return a writable tessdata dir: prefer TESSDATA_PREFIX, then default tessdata;
    if the default is not writable, create a user-cache dir and copy existing packs so tesseract
    finds both old and new."""
    env = os.environ.get("TESSDATA_PREFIX")
    if env and os.path.isdir(env) and _writable(env):
        return env
    default = _default_tessdata_dir()
    if default and _writable(default):
        return default
    cache = os.path.join(os.path.expanduser("~"), ".notes_ocr_tessdata")
    os.makedirs(cache, exist_ok=True)
    if default:
        for f in os.listdir(default):
            if f.endswith(".traineddata") and not os.path.exists(os.path.join(cache, f)):
                try:
                    shutil.copy2(os.path.join(default, f), os.path.join(cache, f))
                except Exception:
                    pass
    return cache


def _download_traineddata(lang: str, tdir: str) -> bool:
    if not tdir or not _writable(tdir):
        return False
    url = f"https://github.com/tesseract-ocr/tessdata/raw/main/{lang}.traineddata"
    dest = os.path.join(tdir, f"{lang}.traineddata")
    try:
        import urllib.request
        urllib.request.urlretrieve(url, dest)
        return os.path.isfile(dest) and os.path.getsize(dest) > 1000
    except Exception as e:
        print(f"⚠️ Failed to download language pack {lang}: {e}", file=sys.stderr)
        return False


def resolve_ocr_langs():
    """Resolve NOTES_OCR_LANG (default chi_sim+deu+eng, covering Chinese/German/English notes);
    try to auto-download missing packs, degrade gracefully if a download fails."""
    tdir = _ensure_tessdata_dir()
    os.environ["TESSDATA_PREFIX"] = tdir  # make tesseract use this dir (includes copied default packs)
    requested = [x.strip() for x in os.environ.get("NOTES_OCR_LANG", "chi_sim+deu+eng").split("+") if x.strip()]
    available = []
    for lang in requested:
        if os.path.isfile(os.path.join(tdir, f"{lang}.traineddata")):
            available.append(lang)
        elif _download_traineddata(lang, tdir):
            available.append(lang)
        else:
            print(f"⚠️ OCR language pack missing and cannot be downloaded, skipping: {lang}", file=sys.stderr)
    if not available:
        raise RuntimeError(
            f"None of the required OCR languages {requested} are available and none could be downloaded. "
            f"Install the language packs manually, or set NOTES_OCR_LANG to an installed language (e.g. eng)."
        )
    return "+".join(available)


def _pdf_pages_to_images(path: Path, indices=None):
    """Scanned PDF → images: prefer PyMuPDF (no external dependency); else pdf2image+poppler.
    Returns a list of (page_no, PIL.Image) or None. `indices` is an optional 1-based
    list of pages to render (None = all)."""
    if PYMUPDF:
        try:
            import io
            from PIL import Image
            import fitz
            doc = fitz.open(str(path))
            images = []
            for idx, page in enumerate(doc, 1):
                if indices and idx not in indices:
                    continue
                pix = page.get_pixmap(dpi=200)
                if pix.alpha:
                    pix = fitz.Pixmap(pix, 0)  # strip alpha channel
                images.append((idx, Image.open(io.BytesIO(pix.tobytes("png")))))
            return images
        except Exception as e:
            print(f"⚠️ PyMuPDF render failed, falling back to poppler: {e}", file=sys.stderr)
    if PDF2IMAGE and PDFTOPPM_BIN:
        from pdf2image import convert_from_path
        poppler_path = os.path.dirname(PDFTOPPM_BIN)
        pages = convert_from_path(str(path), dpi=200, poppler_path=poppler_path)
        out = []
        for idx, img in enumerate(pages, 1):
            if indices and idx not in indices:
                continue
            out.append((idx, img))
        return out
    return None


def _ocr_pdf(path: Path, out_dir: Path, page_indices=None):
    """OCR a PDF. `page_indices` is an optional 1-based list of pages to OCR
    (None = all). Returns dict: page_no -> OCR text.
    Streams page-by-page with PyMuPDF (render one page -> OCR -> release) so
    large image PDFs don't blow up memory; poppler is the fallback path."""
    import pytesseract
    if TESSERACT_BIN:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_BIN
    lang = resolve_ocr_langs()
    result = {}
    rendered = False
    # Preferred: PyMuPDF, rendered one page at a time (low memory even for
    # 100+ page decks)
    if PYMUPDF:
        try:
            import io
            from PIL import Image
            import fitz
            doc = fitz.open(str(path))
            try:
                for idx, page in enumerate(doc, 1):
                    if page_indices and idx not in page_indices:
                        continue
                    pix = page.get_pixmap(dpi=200)
                    if pix.alpha:
                        pix = fitz.Pixmap(pix, 0)  # strip alpha channel
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    result[idx] = pytesseract.image_to_string(img, lang=lang)
                    rendered = True
            finally:
                doc.close()
        except Exception as e:
            print(f"⚠️ PyMuPDF render failed, falling back to poppler: {e}", file=sys.stderr)
    # Fallback: pdf2image + poppler
    if not rendered:
        images = _pdf_pages_to_images(path, indices=page_indices)
        if images:
            for idx, img in images:
                result[idx] = pytesseract.image_to_string(img, lang=lang)
                rendered = True
    if not rendered:
        raise RuntimeError(
            "Scanned PDF needs a rendering engine: install PyMuPDF (`pip install pymupdf`) or poppler (pdftoppm)."
        )
    return result


def _ocr_image(path: Path):
    import pytesseract
    from PIL import Image
    if TESSERACT_BIN:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_BIN
    lang = resolve_ocr_langs()
    img = Image.open(str(path))
    return pytesseract.image_to_string(img, lang=lang)


def parse_pptx(path: Path):
    import pptx
    prs = pptx.Presentation(str(path))
    blocks = []
    for si, slide in enumerate(prs.slides, 1):
        blocks.append(f"\n[Slide {si}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    blocks.append(txt)
            # Bug 8: tables inside shapes were previously dropped — include them
            if shape.has_table:
                tbl = shape.table
                blocks.append("[Table]")
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    blocks.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append(f"(Notes) {notes}")
    return "\n".join(blocks), "python-pptx"


def parse_xlsx(path: Path):
    """Read every sheet of an .xlsx workbook and render each as a Markdown table.
    Uses openpyxl with data_only=True so formula cells resolve to their cached
    values. Multi-line cells are flattened (newlines -> spaces) to keep the
    Markdown table valid; pipe characters are escaped."""
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    blocks = []
    for ws in wb.worksheets:
        blocks.append(f"\n## Sheet: {ws.title}")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            blocks.append("(empty sheet)")
            continue
        ncols = max((len(r) for r in rows), default=0)
        if ncols == 0:
            blocks.append("(no columns)")
            continue
        def fmt(c):
            if c is None:
                return ""
            return str(c).replace("\n", " ").replace("\r", " ").replace("|", "\\|").strip()
        md_rows = []
        for r in rows:
            cells = list(r) + [""] * (ncols - len(r))
            md_rows.append("| " + " | ".join(fmt(c) for c in cells) + " |")
        sep = "| " + " | ".join(["---"] * ncols) + " |"
        blocks.append(md_rows[0] + "\n" + sep + "\n" + "\n".join(md_rows[1:]))
    return "\n".join(blocks), "openpyxl"


def parse_ppt(path: Path, out_dir: Path):
    """Cross-platform .ppt parsing: textutil -> LibreOffice to pptx -> error."""
    if TEXTUTIL:
        return _doc_via_textutil(path, out_dir), "macOS textutil"
    if SOFFICE:
        return _doc_via_soffice(path, out_dir, "pptx"), "LibreOffice"
    raise RuntimeError(
        "Parsing legacy .ppt requires macOS textutil, or LibreOffice (soffice / libreoffice) installed. "
        "On other platforms, convert .ppt to .pptx first."
    )


def parse_image(path: Path, out_dir: Path):
    if not OCR_IMG_OK:
        missing = []
        if not TESSERACT_BIN: missing.append("tesseract")
        if not PYTESSERACT:   missing.append("pytesseract (Python lib)")
        if not PIL_OK:        missing.append("Pillow (PIL)")
        raise RuntimeError(
            f"Image `{path.name}` needs OCR but is missing: {', '.join(missing)}."
            f"\n  Install: macOS `brew install tesseract`;"
            f"Windows `winget install UB-Mannheim.TesseractOCR`;"
            f"Linux `apt install tesseract-ocr`; then `pip install pytesseract pillow`."
        )
    text = _ocr_image(path)
    # optional: also save a PDF (for the "image to PDF" need)
    pdf_path = None
    if IMG2PDF:
        try:
            import img2pdf
            pdf_path = out_dir / f"{path.stem}.pdf"
            with open(pdf_path, "wb") as f:
                f.write(img2pdf.convert(str(path)))
        except Exception:
            pdf_path = None
    return text, "tesseract OCR", pdf_path


def _fmt_time(seconds):
    """Format seconds as M:SS or H:MM:SS."""
    seconds = int(round(seconds))
    h, rem = divmod(seconds, 3600)
    m, s = divmod(rem, 60)
    return f"{h}:{m:02d}:{s:02d}" if h else f"{m:02d}:{s:02d}"


def parse_audio(path: Path, out_dir: Path):
    """Transcribe an audio/video file to text via local Whisper (faster-whisper).
    Returns (markdown_text, engine_label). The model is auto-downloaded from
    HuggingFace on first use and cached; language is auto-detected unless
    NOTES_WHISPER_LANG is set. Audio decoding uses PyAV (bundled ffmpeg), so no
    system ffmpeg install is required."""
    if not py_has("faster_whisper"):
        raise RuntimeError(
            f"Audio `{path.name}` needs faster-whisper (local, free speech-to-text). "
            f"Run `pip install faster-whisper` — it bundles PyAV/ffmpeg, no system ffmpeg needed."
        )
    import faster_whisper
    model = os.environ.get("NOTES_WHISPER_MODEL", "small")
    lang = os.environ.get("NOTES_WHISPER_LANG") or None
    try:
        m = faster_whisper.WhisperModel(model, device="cpu")
    except Exception as e:
        raise RuntimeError(
            f"Failed to load Whisper model '{model}' (auto-downloaded from HuggingFace on first use; "
            f"needs internet on first run): {e}"
        )
    try:
        segments, info = m.transcribe(str(path), language=lang, beam_size=5)
    except Exception as e:
        raise RuntimeError(f"Whisper transcription of `{path.name}` failed: {e}")
    lines = [f"> Transcribed with Whisper model `{model}`"
             + (f", language={lang}" if lang else ", language auto-detected")
             + f"; duration {_fmt_time(info.duration)}; {len(segments)} segments.\n"]
    for seg in segments:
        txt = (seg.text or "").strip()
        if txt:
            lines.append(f"[{_fmt_time(seg.start)}] {txt}")
    return "\n".join(lines), f"Whisper ({model})"


# ---------- main flow ----------
def write_markdown(out_dir: Path, stem: str, raw_text: str, meta: dict, chunk=40000):
    out_dir.mkdir(parents=True, exist_ok=True)
    # Bug 3: remove stale outputs for this stem before (re)writing — prevents
    # orphaned _partN.md files when a re-run produces fewer chunks, and keeps
    # the parsed directory consistent across runs.
    for sf in [out_dir / f"{stem}.md"] + list(out_dir.glob(f"{stem}_part*.md")):
        _safe_unlink(sf)
    header = (
        f"# {meta['title']}\n\n"
        f"> Source: `{meta['source']}`\n"
        f"> Format: {meta['format']}\n"
        f"> Parser: {meta['engine']}\n\n---\n\n"
    )
    text = raw_text.strip()
    if not text:
        text = "(No text content could be extracted from this file)"
    # split oversized files into chunks
    if len(text) <= chunk:
        (out_dir / f"{stem}.md").write_text(header + text, encoding="utf-8")
        return [f"{stem}.md"]
    parts = []
    for i in range(0, len(text), chunk):
        part = text[i:i + chunk]
        fn = f"{stem}_part{i // chunk + 1}.md"
        (out_dir / fn).write_text(header + part, encoding="utf-8")
        parts.append(fn)
    return parts


def report_tools():
    rows = [
        ("python-docx", PYTHON_DOCX),
        ("pdfplumber", PDFPLUMBER),
        ("python-pptx", PYTHON_PPTX),
        ("openpyxl", OPENPYXL),
        ("img2pdf", IMG2PDF),
        ("pytesseract", PYTESSERACT),
        ("Pillow(PIL)", PIL_OK),
        ("PyMuPDF(fitz)", PYMUPDF),
        ("pdf2image", PDF2IMAGE),
        ("tesseract binary", TESSERACT_BIN),
        ("pdftoppm(poppler)", PDFTOPPM_BIN),
        ("macOS textutil", TEXTUTIL),
        ("Windows Word COM (pywin32)", WIN32COM),
        ("LibreOffice (soffice)", SOFFICE),
        ("antiword", ANTIWORD),
        ("PyAV / ffmpeg (via faster-whisper)", py_has("av")),
        ("faster-whisper", WHISPER_OK),
    ]
    print("=== Tool availability self-check ===")
    for name, ok in rows:
        print(f"  [{'✔' if ok else '✘'}] {name}")
    ocr_state = "✔ available" if OCR_PDF_OK else ("partial (image OCR only)" if OCR_IMG_OK else "✘ unavailable")
    print(f"  >> Overall OCR status: {ocr_state}")
    audio_state = "✔ available" if WHISPER_OK else "✘ unavailable (audio/recording files will be skipped)"
    print(f"  >> Speech-to-text (Whisper) status: {audio_state}\n")


def _log_done(log, name, p):
    """Append a processed-file record (name ⇥ mtime ⇥ size) to processed.log."""
    try:
        m = p.stat()
        with log.open("a", encoding="utf-8") as f:
            f.write(f"{name}\t{int(m.st_mtime)}\t{m.st_size}\n")
    except Exception:
        pass


def _unchanged(p, rec):
    """True if the file's current mtime+size match the recorded processed.log entry."""
    try:
        m = p.stat()
        return rec is not None and rec == (int(m.st_mtime), m.st_size)
    except Exception:
        return False


def _dispatch_and_write(p, out, out_stem=None):
    """Parse a single regular (non-archive) file and write its Markdown.
    Returns True on success, False on failure (engine missing, parse error).
    Does NOT update processed.log — the caller owns log writes."""
    ext = p.suffix.lower()
    try:
        if ext in IMAGE_EXTS:
            res = parse_image(p, out)
            if len(res) == 3:
                text, engine, pdf_path = res
                extra = f" (saved PDF: {pdf_path.name})" if pdf_path else ""
            else:
                text, engine = res
                extra = ""
        elif ext == ".docx":
            text, engine = parse_docx(p)
        elif ext == ".doc":
            text, engine = parse_doc(p, out)
        elif ext == ".pdf":
            text, engine = parse_pdf(p, out)
        elif ext == ".pptx":
            text, engine = parse_pptx(p)
        elif ext == ".ppt":
            text, engine = parse_ppt(p, out)
        elif ext == ".xlsx":
            text, engine = parse_xlsx(p)
        elif ext in AUDIO_EXTS:
            text, engine = parse_audio(p, out)
        else:
            return False
        out_stem = out_stem or f"{p.stem}{ext}"  # Bug 1: embed ext so same-name/diff-format files don't collide
        parts = write_markdown(
            out, out_stem, text,
            {"title": p.name, "source": str(p), "format": ext, "engine": engine}
        )
        label = extra if ext in IMAGE_EXTS else ""
        print(f"✅ {p.name} -> {', '.join(parts)} [{engine}]{label}")
        return True
    except Exception as e:
        print(f"⚠️ {p.name} skipped: {e}", file=sys.stderr)
        return False


def handle_archive(path, out):
    """Extract a .zip/.rar archive to a temp dir and recurse into its contents.
    Each supported inner file is parsed and written as its own Markdown; nested
    archives are extracted too. Returns the list of inner files successfully
    parsed. Does NOT write processed.log (the caller logs the archive itself)."""
    import zipfile
    import tempfile
    import shutil
    archive_stem = path.stem
    parsed_names = []
    tmp = Path(tempfile.mkdtemp(prefix=f"notes_arch_{archive_stem}_"))
    try:
        try:
            if path.suffix.lower() == ".zip":
                with zipfile.ZipFile(str(path)) as z:
                    z.extractall(str(tmp))
            elif path.suffix.lower() == ".rar":
                try:
                    import rarfile
                except ImportError:
                    print(f"⚠️ `{path.name}` is a .rar but `rarfile` is not installed; "
                          f"run `pip install rarfile` then install a rar tool "
                          f"(macOS `brew install unar`; Linux `apt install unrar`).", file=sys.stderr)
                    return parsed_names
                try:
                    with rarfile.RarFile(str(path)) as rf:
                        rf.extractall(str(tmp))
                except Exception as e:
                    print(f"⚠️ Failed to extract .rar `{path.name}` (need unrar/unar/bsdtar): {e}",
                          file=sys.stderr)
                    return parsed_names
            else:
                return parsed_names
        except Exception as e:
            print(f"⚠️ Failed to extract archive `{path.name}`: {e}", file=sys.stderr)
            return parsed_names

        # rglob over extracted contents; ARCHIVE_EXTS included so nested archives recurse.
        for f in sorted(tmp.rglob("*")):
            if not f.is_file():
                continue
            fext = f.suffix.lower()
            if fext not in SUPPORTED:
                continue
            rel = f.relative_to(tmp)
            slug = str(rel.with_suffix("")).replace(os.sep, "__")
            out_stem = f"{archive_stem}__{slug}{fext}"
            if fext in ARCHIVE_EXTS:
                # nested archive: recurse (inner out-stem already prefixed)
                parsed_names.extend(handle_archive(f, out))
            else:
                if _dispatch_and_write(f, out, out_stem=out_stem):
                    parsed_names.append(f"{archive_stem}/{rel}")
        return parsed_names
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def main():
    ap = argparse.ArgumentParser(description="Unified free parser: study notes → Markdown (cross-platform)")
    ap.add_argument("--source", required=True, help="Notes source folder")
    ap.add_argument("--out", default=None, help="Parsed output directory (default <source>/../parsed)")
    ap.add_argument("--log", default=None, help="processed.log path (default <out>/processed.log)")
    ap.add_argument("--force", action="store_true", help="Ignore processed.log, force re-parse everything")
    args = ap.parse_args()

    src = Path(args.source).expanduser().resolve()
    if not src.is_dir():
        print(f"❌ Source folder does not exist: {src}", file=sys.stderr)
        sys.exit(1)

    out = Path(args.out).expanduser().resolve() if args.out else (src.parent / "parsed")
    out.mkdir(parents=True, exist_ok=True)
    log = Path(args.log).expanduser().resolve() if args.log else (out / "processed.log")

    report_tools()

    done = {}  # name -> (mtime, size)  — Bug 2: re-parse when a file changes
    if log.exists() and not args.force:
        for line in log.read_text(encoding="utf-8", errors="ignore").splitlines():
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            parts = line.split("\t")
            if len(parts) >= 3:
                try:
                    done[parts[0]] = (int(parts[1]), int(parts[2]))
                except ValueError:
                    done[parts[0]] = None
            else:
                done[parts[0]] = None  # legacy one-column log

    # Bug 10: descend into sub-directories too; also count unsupported files for exit-code logic
    all_files = sorted([p for p in src.rglob("*") if p.is_file()])
    files = [p for p in all_files if p.suffix.lower() in SUPPORTED]
    unsupported_count = sum(1 for p in all_files if p.suffix.lower() not in SUPPORTED)
    print(f"Source folder: {src}")
    print(f"Output folder: {out}")
    print(f"Files to process: {len(files)}\n")

    parsed, skipped = [], []
    for p in files:
        ext = p.suffix.lower()
        if not args.force and p.name in done and _unchanged(p, done[p.name]):
            continue
        if ext in ARCHIVE_EXTS:
            # Auto-extract the archive and recurse into its contents (zip/rar).
            inner = handle_archive(p, out)
            if inner:
                parsed.extend(inner)
                _log_done(log, p.name, p)
            else:
                skipped.append((p.name, "archive empty, extraction failed, or no supported files inside"))
            continue
        ok = _dispatch_and_write(p, out)
        if ok:
            parsed.append(p.name)
            _log_done(log, p.name, p)
        else:
            skipped.append((p.name, "parse failed / engine missing"))

    print(f"\n=== Done ===")
    print(f"Parsed this run: {len(parsed)}; skipped: {len(skipped)}")
    if skipped:
        print("Skipped files (re-run after fixing the environment to auto-retry):")
        for n, r in skipped:
            print(f"  - {n}: {r}")
    # Bug 9: exit code — 2 if nothing was parsed but files existed (all skipped,
    # unsupported, or engine-missing), else 0
    return 2 if (not parsed and (skipped or unsupported_count)) else 0


if __name__ == "__main__":
    sys.exit(main() or 0)
